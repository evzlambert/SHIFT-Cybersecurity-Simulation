"""
Build SHIFT Cybersecurity Breach Simulation PowerPoint.
Generates a 35-slide .pptx matching the instructor app with speaker notes.

Usage: python3 build_pptx.py
Output: SHIFT-Cybersecurity-Simulation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------
BG_DARK = RGBColor(0x0A, 0x0A, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
RED = RGBColor(0xE9, 0x45, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xF0)
MID_GRAY = RGBColor(0xA0, 0xA0, 0xC0)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)
GREEN_BG = RGBColor(0x2E, 0x7D, 0x32)
PURPLE_BG = RGBColor(0x4A, 0x14, 0x8C)
PURPLE_LIGHT = RGBColor(0xCE, 0x93, 0xD8)
ORANGE = RGBColor(0xFF, 0xA7, 0x26)
ORANGE_DARK = RGBColor(0xE6, 0x51, 0x00)
BLUE = RGBColor(0x42, 0xA5, 0xF5)
BLUE_LIGHT = RGBColor(0xBB, 0xDE, 0xFB)
RED_BANNER = RGBColor(0xB7, 0x1C, 0x1C)
SYS_GREEN = RGBColor(0x2E, 0x7D, 0x32)
SYS_RED = RGBColor(0xC6, 0x28, 0x28)
SYS_ORANGE = RGBColor(0xE6, 0x51, 0x00)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Narrations (from generate_audio.py)
# ---------------------------------------------------------------------------
NARRATIONS = {
    0: "Welcome to the SHIFT Cybersecurity Breach Simulation. You are the senior leadership team of Centre Hospitalier Sainte-Claire, a 350-bed regional hospital in southeastern France. 1,200 staff. 287 inpatients. The only Level 1 trauma center and neonatal ICU in the region. It is Monday. 6:12 AM. All systems are operational. That is about to change.",
    1: "6:12 AM. A night-shift nurse in the ICU reports that the electronic health record is frozen on three workstations. 6:15. A second nurse reports the same issue in the maternity unit. She cannot access patient birth plans. 6:18. The IT help desk receives twelve simultaneous calls. The lab system is showing error messages. 6:22. The pharmacy technician reports that automated dispensing cabinets are unresponsive. 6:24. Radiology PACS displays a ransom message. Your files are encrypted. Contact us for the decryption key. You have 72 hours.",
    2: "6:26. The night IT administrator attempts to restart the servers. The same ransom message appears on four of six core servers. 6:30. The pharmacy dispensing system is confirmed offline. Medication cabinets are locked hospital-wide. 6:33. The patient portal goes dark. 47 telehealth appointments were scheduled for today. 6:35. The scheduling system is unresponsive. 142 outpatient appointments cannot be confirmed. 6:38. Billing and revenue cycle, encrypted. No claims can be submitted. 6:40. IT Security confirms: ransomware, a LockBit variant, has encrypted servers across the EHR, lab system, PACS, pharmacy, and billing. 6:42. Forensic log review reveals that the initial access occurred eleven days ago, through a phishing email sent to a radiology technician.",
    3: "6:45. Email servers are compromised. Internal communication is now limited to phone, in-person, and personal devices. 6:48. The nurse call system is functioning but degraded. Some units report delays. 6:50. Badge access is intermittent. The pharmacy and server room cannot be secured electronically. 6:55. The 7 AM shift change begins. More than 400 day-shift staff are arriving to a hospital in crisis. Most have no information. 7:00. The Hospital Incident Command System is activated. 7:05. ANSSI, France's national cybersecurity agency, has been contacted. They estimate a forensics team can arrive in 6 to 8 hours. All major clinical systems are down.",
    4: "Pause. This is your first team discussion. You are the hospital's crisis leadership team. It is 7:05 AM. Discuss with your team. First: What is your first action? Do you activate downtime procedures or isolate the network? Second: How do you communicate to 1,200 staff with no email? Third: 15 ICU patients are on smart pump drips. 8 women are in labor. 12 surgeries are scheduled. What is your patient safety strategy? Fourth: Who do you call first? Police, insurance, the regulator, or the board? Discuss on your feet. Then return to your device to log your decisions.",
    5: "Downtime operations. Monday 8 AM through Tuesday 6 AM. Hours 2 through 24. All clinical systems are operating manually. Every medication administration. Every patient identification check. Every lab order. Every clinical note. On paper. 287 inpatients. 42 in the ICU. 8 in maternity. 14 patients named Martin. Zero electronic systems online.",
    6: "Inject event. Hour 4. Near-miss medication event. A nurse on the medical-surgical floor reports a near-miss. Without barcode scanning, a patient with a similar name, Martin, Jean-Pierre, versus Martin, Jean-Paul, was almost given the wrong medication. The manual five-rights check caught the error. But staff are shaken, and they are requesting additional safety protocols.",
    7: "Pause. Team discussion. A near-miss just occurred. Your staff are scared. From your role's perspective, what is the immediate response? How do you prevent the next one from being an actual error? Discuss, then log your decision in the student app.",
    8: "Inject event. Hour 8. Ambulance diversion request. The Emergency Department is overwhelmed. Without electronic triage, patient tracking, or access to prior medical histories, ED wait times have tripled. The ED Medical Director is requesting ambulance diversion. However, the nearest Level 1 trauma center is 45 minutes away. Your hospital has the only neonatal ICU in the region. And right now, a pregnant woman at 36 weeks is in an ambulance headed toward you.",
    9: "Pause. Team discussion. An ambulance with a pregnant patient is en route. Your ED is overwhelmed. Do you approve full diversion? Partial diversion? Or deny it entirely? And what do you tell the paramedics in that ambulance right now? Discuss, then log your decision.",
    10: "Inject event. Hour 13. The 7 PM shift handoff. The evening shift is arriving. Day-shift staff have been managing the crisis for thirteen hours. Evening staff have received minimal information. You must hand off: Current status of 287 inpatients, all documented on paper. Active downtime procedures and which workflows are in place. Outstanding issues, including the medication near-miss and the ambulance diversion status. The location of all paper records generated today. And the emotional state of your day-shift staff. Several nurses are in tears.",
    11: "Pause. Team discussion. Evening staff are arriving. They know nothing about paper processes. How do you structure the handoff? Do you keep your exhausted day-shift staff longer, or let them go? Discuss, then log your decision.",
    12: "Inject event. Hour 18. Paper records crisis. After 18 hours of downtime, an estimated 3,200 pages of handwritten records have accumulated across the hospital. A quality audit of 50 paper records reveals: 34 percent are missing the patient's date of birth. 22 percent have no allergies documented. 18 percent have illegible handwriting in at least one section. 11 percent cannot be matched to a specific patient with certainty. And 4 records appear to be duplicates for the same patient.",
    13: "Pause. Team discussion. 3,200 pages of paper. 11 percent unmatched to any patient. Illegible handwriting. Exhausted staff. What is your data management strategy? And how will you handle back-entry when systems return? Discuss, then log your decision. After this, we take a break.",
    14: "Break. 20 minutes. The crisis does not pause, but you can. When you return, it will be Hour 24. The media are calling. Families are arriving. Staff are reaching their breaking point. The pressure is about to intensify.",
    15: "Day 2. Escalation and stakeholder pressure. Tuesday. Hours 24 through 48. 24 hours of downtime. 3,200 paper records. 23 surgeries canceled. 14 patients diverted. An estimated 240,000 euros in lost revenue. And now, the external pressure begins.",
    16: "Inject event. Hour 26. The story breaks. France 3 Auvergne-Rhone-Alpes, noon broadcast. Quote: Centre Hospitalier Sainte-Claire has been operating without its computer systems for over 24 hours following what sources describe as a major cyberattack. Ambulances have been diverted from the hospital's emergency department since yesterday morning. A hospital spokesperson confirmed a technology disruption but declined to provide details. Patient advocacy groups are demanding transparency. The regional health authority says it is monitoring the situation. End quote. The phone is ringing. Media, patients, families. How do you respond?",
    17: "Pause. Team discussion. The story is on television. Your phone lines are overwhelmed. Do you hold a press conference? Issue a written statement? Go fully transparent? Who is the spokesperson? What is the one key message? Discuss, then log your decision.",
    18: "Inject event. Hour 28. A family demands answers. A man approaches the hospital information desk. He is visibly upset. His mother, 82 years old, was admitted three days ago for pneumonia. He has been unable to reach anyone by phone. The patient portal is down. He drove 90 minutes to get here. He says: I cannot access my mother's records. I don't know what medications she's on. The nurse told me they are using paper. Paper! What year is this? I want to know if my mother's personal information has been stolen. I want to speak to whoever is in charge. And if I don't get answers, I am calling my lawyer and the press.",
    19: "Pause. Team discussion. This man represents hundreds of families with the same fears. Who meets with him? What information can you share? He is threatening legal action and media contact. How do you de-escalate? Discuss, then log your decision.",
    20: "Inject event. Hour 32. Staff morale crisis. Voicemail from Marie-Claire, charge nurse, Unit 3B. Quote: This is Marie-Claire on 3B. I need to tell you that my nurses are done. We have been working 14-hour shifts on paper for two days. Three of my nurses called in sick today. I think they're not actually sick. They just can't take it anymore. The ones who are here are making mistakes. Small ones, so far. But I'm scared. We had another near-miss. Wrong dose calculated by hand because we don't have the EHR doing the math for us. I need more staff. I need someone to acknowledge what we're going through. And I need to know when this is going to end. Please call me back. End quote.",
    21: "Pause. Team discussion. Marie-Claire is one of your best charge nurses. If she breaks, the unit breaks. Do you request mutual aid? Reduce your patient census? Deploy crisis counseling? Shorten shifts? What does each of those cost you? Discuss, then log your decision.",
    22: "Inject event. Hour 34. Social media firestorm. Multiple posts are circulating. HealthWatch FR, with 12,000 followers, writes: Thread. My aunt is a nurse at Sainte-Claire. She says they've been writing everything by hand for two days. No access to patient histories. No medication barcodes. She's terrified of making an error. This is a patient safety emergency and no one is talking about it. Patient Rights FR, with 34,000 followers, writes: If Sainte-Claire has had 180,000 patient records stolen, where is the CNIL notification? GDPR requires transparency. Silence is not acceptable. And a verified physician, not affiliated with the hospital, writes: The Sainte-Claire situation is exactly what we've been warning about. French hospitals are chronically under-invested in cybersecurity. This will happen again. And next time, patients may die.",
    23: "Pause. Team discussion. A nurse's family member posted. Advocacy groups are demanding answers. A doctor is predicting patient harm. Do you post an official response on social media? Ignore it? Ask the employee to have the post removed? Hire a crisis communications firm? Discuss, then log your decision.",
    24: "Inject event. Hour 36. Inter-department conflict. Scene: the incident command center. The Chief of Surgery and the HIM Director are in a heated argument. Surgery wants to resume elective procedures tomorrow morning. Patients are at medical risk from delay. The HIM Director insists that without electronic surgical checklists, specimen labeling systems, and proper consent documentation workflows, resuming surgery on paper is an unacceptable risk. The Chief of Surgery says: We operated for decades before computers. I am not letting a computer problem endanger my patients by delaying their care. The HIM Director responds: And we had higher error rates for decades before computers. We are not going backward. Every mislabeled specimen is a potential catastrophe. Both are right. Both are wrong. The Incident Commander must decide.",
    25: "Pause. Team discussion. Surgery versus HIM. Patient harm from delay versus patient harm from error. Do you resume elective surgery on paper? Delay until systems return? Find a compromise? Convene a governance meeting? Discuss, then log your decision.",
    26: "Day 3. Critical decisions. Wednesday. Hours 48 through 72. The ransom deadline is approaching. The GDPR clock is ticking. 48 hours of downtime. 6,400 paper records. 38 surgeries canceled. 480,000 euros in estimated revenue loss. 89 patient complaints. The hardest decisions of this crisis are ahead.",
    27: "Inject event. Hour 50. The ransom demand. The attackers have made formal contact through an encrypted channel. They demand 2.5 million euros in cryptocurrency within 48 hours. They claim to have stolen 180,000 patient records, including diagnoses, medications, national identification numbers, and in some cases, mental health and HIV status records. They provide a sample of 500 records as proof. Your IT team confirms: these are genuine. They threaten to publish all data on the dark web if payment is not received. Your cyber insurance policy covers 1 million euros. ANSSI officially advises against paying. Under French law, the Loi LOPMI of 2023, a police report must be filed within 72 hours of any ransom payment to be eligible for insurance reimbursement. 2.5 million euros. 180,000 patients. Mental health records. HIV status. What do you recommend to the board?",
    28: "Pause. This is the hardest decision in the simulation. Take your time with this one. Do you pay? Negotiate? Refuse? Escalate to the government? Consider the ethical dimensions. Who bears the harm in each scenario? Is it ethical to use hospital funds to pay criminals? Is it ethical not to pay, if patient data will be published? Discuss thoroughly, then log your decision and your reasoning.",
    29: "Inject event. Hour 58. GDPR 72-hour notification deadline. Under GDPR Article 33, you must notify the CNIL, France's data protection authority, within 72 hours of becoming aware of a personal data breach. The clock started when the breach was confirmed at 6:40 Monday. You are now at hour 58. You have 14 hours remaining. The forensics team has established that the attackers had network access for 11 days before encryption. Large data transfers to external servers have been confirmed. But the exact number of records exfiltrated is still unknown. The 500-record sample provided by the attackers has been verified as genuine. The breach likely includes patient demographics, diagnoses, medications, lab results, and for some patients, mental health records. 14 hours remain on the clock. Incomplete information. 180,000 potentially affected patients.",
    30: "Pause. Team discussion. 14 hours on the clock. Incomplete information. Do you file a preliminary notification now and supplement later? Wait for the deadline and file the most complete report possible? And critically: how do you notify 180,000 patients that their health data may have been stolen? Consider: literacy levels. Elderly patients. Patients with mental health records who may face stigma. Domestic violence survivors whose addresses are now potentially exposed. Discuss, then log your decision.",
    31: "Inject event. Hour 64. System restoration strategy. The IT team presents four options. Option A: Restore from a 48-hour-old backup. The EHR can be operational in 12 to 18 hours. But you lose 48 hours of pre-breach data, plus everything documented on paper during downtime. The backup is verified clean. Back-entry will take 3 to 4 weeks. Option B: Attempt to decrypt and clean the current systems. This could take 5 to 10 additional days. It preserves all pre-breach data, but carries a 15 to 20 percent risk of reinfection. The forensics team recommends against this option. Option C: A new, clean EHR instance on new hardware. Operational in 24 to 36 hours. Zero reinfection risk. But it requires manual entry of all active patient data. Historical data would only be available later. Option D: A hybrid phased approach. Phase 1: stand up a clean instance for active inpatients in 24 hours. Phase 2: restore the backup in a parallel environment over one week. Phase 3: merge, reconcile, and back-enter paper records over two to four weeks. The most complex option, but the most balanced. Speed versus data preservation versus security. Who decides?",
    32: "Pause. Final team discussion. Four options. Each has costs. Each has risks. Which restoration path do you choose? Is this a technical decision? A clinical decision? A financial decision? A data integrity decision? Who has the final say? Discuss, then log your final decision.",
    33: "Recovery planning. Thursday. Hour 72 and beyond. Partial systems restoration is in progress. Some systems are coming back online. But the work is far from over. The systems will come back in days. The trust takes months. The lessons last careers. In your student app, complete the recovery planning phase. Prioritize the back-entry of paper records. Plan your trust recovery strategy. Assess your hospital's preparedness. Write your after-action summary.",
    34: "End of simulation. Thank you for your work over the past three hours. Export your team reports from the student app. Prepare to share your most difficult decision and what you learned. The two-hour debrief will follow.",
}

# ---------------------------------------------------------------------------
# Systems status configurations
# ---------------------------------------------------------------------------
SYSTEM_NAMES = ["EHR (DPI)", "Lab (LIS)", "PACS", "Pharmacy", "Patient Portal",
                "Email", "Scheduling", "Billing", "Nurse Call", "Badge Access"]

# Status: "UP", "DOWN", "DEGRADED"
ALL_UP = ["UP"] * 10
ALL_DOWN = ["DOWN"] * 8 + ["DEGRADED", "DEGRADED"]
DAY2 = ["DOWN"] * 8 + ["UP", "UP"]
DAY3 = ["DOWN", "DOWN", "DOWN", "DOWN", "DOWN", "DEGRADED", "DOWN", "DOWN", "UP", "UP"]
RECOVERY = ["DEGRADED", "DEGRADED", "DOWN", "DEGRADED", "DOWN", "UP", "DEGRADED", "DOWN", "UP", "UP"]

STATUS_COLORS = {"UP": SYS_GREEN, "DOWN": SYS_RED, "DEGRADED": SYS_ORANGE}

# ---------------------------------------------------------------------------
# Slide definitions
# ---------------------------------------------------------------------------
# Each: (type, title, subtitle, body_lines, extra_data)
# extra_data can be: systems list, metrics list, events list, etc.

SLIDES = [
    # 0: Title
    {
        "type": "title",
        "title": "SHIFT CYBERSECURITY BREACH SIMULATION",
        "subtitle": "Centre Hospitalier Sainte-Claire  |  Saint-Etienne, France",
        "body": [
            "A 350-bed regional hospital. 1,200 staff. 287 inpatients.",
            "The only Level 1 trauma center and neonatal ICU in the region.",
            "",
            "IT IS MONDAY, 06:12 AM.",
            "",
            "All systems operational.",
        ],
        "systems": ALL_UP,
    },
    # 1: Breach begins
    {
        "type": "breach",
        "title": "THE BREACH BEGINS",
        "subtitle": "Monday -- 06:12 to 06:24",
        "events": [
            ("06:12", 'Night-shift nurse reports EHR is "frozen" on 3 workstations in the ICU.'),
            ("06:15", "Second nurse reports the same in maternity. Cannot access patient birth plans."),
            ("06:18", "IT help desk receives 12 simultaneous calls. Lab system showing errors."),
            ("06:22", "Pharmacy technician: automated dispensing cabinets unresponsive."),
            ("06:24", 'Radiology PACS displays a ransom message:\n"Your files are encrypted. Contact us for the decryption key. You have 72 hours."'),
        ],
    },
    # 2: Systems falling
    {
        "type": "breach",
        "title": "SYSTEMS FALLING",
        "subtitle": "Monday -- 06:26 to 06:42",
        "events": [
            ("06:26", "Night IT admin restarts servers -- same ransom message on 4 of 6 core servers."),
            ("06:30", "Pharmacy dispensing confirmed offline. Medication cabinets locked hospital-wide."),
            ("06:33", "Patient portal goes dark. 47 telehealth appointments scheduled today."),
            ("06:35", "Scheduling unresponsive. 142 outpatient appointments unconfirmable."),
            ("06:38", "Billing and revenue cycle encrypted. No claims can be submitted."),
            ("06:40", "IT confirms: LockBit ransomware across EHR, LIS, PACS, pharmacy, billing."),
            ("06:42", "Forensics: initial access was 11 days ago via phishing email to radiology tech."),
        ],
    },
    # 3: Full attack
    {
        "type": "breach",
        "title": "FULL-SCALE ATTACK CONFIRMED",
        "subtitle": "Monday -- 06:45 to 07:05",
        "events": [
            ("06:45", "Email servers compromised. Communication limited to phone and in-person."),
            ("06:48", "Nurse call system degraded. Some units report delays."),
            ("06:50", "Badge access intermittent. Pharmacy and server room unsecured."),
            ("06:55", "400+ day-shift staff arriving. Most have no information."),
            ("07:00", "Hospital Incident Command System (HICS) activated."),
            ("07:05", "ANSSI contacted. Forensics team: 6-8 hours away."),
        ],
        "banner": "ALL MAJOR CLINICAL SYSTEMS ARE DOWN",
        "systems": ALL_DOWN,
    },
    # 4: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "You are the hospital's crisis leadership team. It is 07:05 AM.",
        "prompts": [
            "1. What is your FIRST action? Downtime procedures or network isolation?",
            "2. How do you communicate to 1,200 staff with no email?",
            "3. 15 ICU patients on smart pump drips. 8 women in labor. 12 surgeries scheduled.\n   What is your patient safety strategy?",
            "4. Who do you call first? Police, insurance, regulator, or the board?",
        ],
        "footer": "Discuss on your feet. Then return to your device to log your decisions.",
    },
    # 5: Downtime
    {
        "type": "downtime",
        "title": "DOWNTIME OPERATIONS",
        "subtitle": "Monday 08:00 -- Tuesday 06:00  |  Hour 2 to 24",
        "banner": "ALL CLINICAL SYSTEMS OPERATING MANUALLY",
        "body": [
            "Every medication. Every patient ID check. Every lab order. Every note.",
            "On paper.",
        ],
        "systems": ALL_DOWN,
        "metrics": [("287", "Inpatients"), ("42", "ICU Patients"), ("8", "Maternity"),
                    ("14", "Patients Named Martin"), ("0", "Systems Online")],
    },
    # 6: Near-miss
    {
        "type": "inject",
        "title": "INJECT: NEAR-MISS MEDICATION EVENT",
        "subtitle": "Hour 4",
        "events": [
            ("", "A nurse reports a near-miss: without barcode scanning,\nMARTIN, Jean-Pierre vs. MARTIN, Jean-Paul --\nalmost given the wrong medication."),
            ("", "The manual five-rights check caught it.\nBut staff are shaken and requesting additional safety protocols."),
        ],
        "decision": "How does your team prevent actual medication errors during downtime?",
    },
    # 7: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "A near-miss just occurred. Your staff are scared.",
        "prompts": [
            "From your role's perspective: What is the immediate response?",
            "How do you prevent the next one from being an actual error?",
        ],
        "footer": "Discuss, then log your decision in the student app.",
    },
    # 8: Ambulance
    {
        "type": "inject",
        "title": "INJECT: AMBULANCE DIVERSION REQUEST",
        "subtitle": "Hour 8",
        "events": [
            ("", "The ED is overwhelmed. Without electronic triage or patient tracking,\nwait times have tripled. The ED Medical Director requests diversion."),
            ("", "Nearest Level 1 trauma center: 45 minutes away.\nOnly neonatal ICU in the region: yours.\nA pregnant woman at 36 weeks is in an ambulance en route to you right now."),
        ],
        "decision": "Do you approve diversion? What about the woman in the ambulance?",
    },
    # 9: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "An ambulance with a pregnant patient is en route. Your ED is overwhelmed.",
        "prompts": [
            "Full diversion? Partial? Deny entirely?",
            "What do you tell the paramedics right now?",
        ],
        "footer": "Discuss, then log your decision.",
    },
    # 10: Shift handoff
    {
        "type": "inject",
        "title": "INJECT: THE 19:00 SHIFT HANDOFF",
        "subtitle": "Hour 13",
        "events": [
            ("", "Evening shift arriving. Day shift: 13 hours of crisis management.\nEvening staff have minimal information."),
            ("", "You must hand off:\n"
                 "  - 287 inpatients -- all on paper\n"
                 "  - Active downtime procedures\n"
                 "  - Outstanding issues (near-miss, diversion status)\n"
                 "  - Location of all paper records\n"
                 "  - Emotional state of day-shift staff -- several nurses in tears"),
        ],
        "decision": "How do you hand off a paper-based hospital to staff who've never used paper?",
    },
    # 11: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "Evening staff arriving. They know nothing about paper processes.",
        "prompts": [
            "How do you structure the handoff?",
            "Keep exhausted day-shift longer, or let them go?",
        ],
        "footer": "Discuss, then log your decision.",
    },
    # 12: Paper crisis
    {
        "type": "inject",
        "title": "INJECT: PAPER RECORDS CRISIS",
        "subtitle": "Hour 18",
        "events": [
            ("", "3,200 pages of handwritten records accumulated.\nQuality audit of 50 records:"),
        ],
        "metrics": [("34%", "Missing DOB"), ("22%", "No Allergies"), ("18%", "Illegible"),
                    ("11%", "Can't Match\nto Patient"), ("4", "Duplicates")],
        "decision": "How do you manage the backlog while maintaining data integrity?",
    },
    # 13: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "3,200 pages. 11% unmatched. Illegible handwriting. Exhausted staff.",
        "prompts": [
            "What is your data management strategy?",
            "How will you handle back-entry when systems return?",
        ],
        "footer": "Discuss, then log your decision. After this, we take a break.",
    },
    # 14: Break
    {
        "type": "break",
        "title": "BREAK",
        "body": [
            "20 minutes. The crisis does not pause, but you can.",
            "",
            "When you return, it will be Hour 24.",
            "The media are calling. Families are arriving. Staff are breaking.",
            "The pressure is about to intensify.",
        ],
    },
    # 15: Day 2
    {
        "type": "escalation",
        "title": "ESCALATION & STAKEHOLDER PRESSURE",
        "subtitle": "Tuesday -- Hour 24 to 48  |  Day 2",
        "banner": "24 HOURS OF DOWNTIME -- EXTERNAL PRESSURE MOUNTING",
        "systems": DAY2,
        "metrics": [("24h", "Downtime"), ("3,200", "Paper Records"), ("23", "Surgeries Canceled"),
                    ("14", "Diverted"), ("\u20ac240K", "Revenue Lost")],
    },
    # 16: Media
    {
        "type": "inject",
        "title": "THE STORY BREAKS",
        "subtitle": "Hour 26",
        "news": {
            "source": "BREAKING -- France 3 Auvergne-Rhone-Alpes -- 12:00 Noon",
            "text": ('"Centre Hospitalier Sainte-Claire has been operating without its computer systems '
                     'for over 24 hours following what sources describe as a \'major cyberattack.\' '
                     'Ambulances have been diverted. A hospital spokesperson confirmed a \'technology '
                     'disruption\' but declined to provide details. Patient advocacy groups demand '
                     'transparency. The ARS says it is \'monitoring the situation.\'"'),
        },
        "decision": "The phone is ringing. Media, patients, families. How do you respond?",
    },
    # 17: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "The story is on television. Your phone lines are overwhelmed.",
        "prompts": [
            "Press conference? Written statement? Full transparency?",
            "Who is the spokesperson? What is the ONE key message?",
        ],
        "footer": "Discuss, then log your decision.",
    },
    # 18: Family
    {
        "type": "inject",
        "title": "A FAMILY DEMANDS ANSWERS",
        "subtitle": "Hour 28",
        "scene": {
            "label": "Scene: Hospital Main Lobby",
            "text": ('A man approaches the information desk. Visibly upset. His mother, 82, '
                     'admitted three days ago for pneumonia. Unable to reach anyone by phone. '
                     'Portal down. Drove 90 minutes.\n\n'
                     '"I can\'t access my mother\'s records. I don\'t know what medications she\'s on. '
                     'The nurse told me they\'re using PAPER. Paper! What year is this? '
                     'I want to know if my mother\'s personal information has been stolen. '
                     'I want to speak to whoever is in charge. And if I don\'t get answers, '
                     'I\'m calling my lawyer and the press."'),
        },
        "decision": "Who responds? What do you say? What can't you say?",
    },
    # 19: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "This man represents hundreds of families with the same fears.",
        "prompts": [
            "Who meets with him? What can you share?",
            "He's threatening legal action and media. How do you de-escalate?",
        ],
        "footer": "Discuss, then log your decision.",
    },
    # 20: Staff morale
    {
        "type": "inject",
        "title": "STAFF MORALE CRISIS",
        "subtitle": "Hour 32",
        "voicemail": {
            "label": "Voicemail -- Marie-Claire, Charge Nurse, Unit 3B",
            "text": ('"This is Marie-Claire on 3B. I need to tell you that my nurses are done. '
                     "We've been working 14-hour shifts on paper for two days. Three of my nurses "
                     "called in sick today -- I think they're not actually sick, they just can't take "
                     "it anymore. The ones who are here are making mistakes. Small ones so far. But "
                     "I'm scared. We had another near-miss -- wrong dose calculated by hand because we "
                     "don't have the EHR doing the math for us. I need more staff, I need someone to "
                     "acknowledge what we're going through, and I need to know when this is going to end. "
                     'Please call me back."'),
        },
        "decision": "Staff burnout is becoming a patient safety issue. What do you do?",
    },
    # 21: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "Marie-Claire is one of your best charge nurses. If she breaks, the unit breaks.",
        "prompts": [
            "Mutual aid? Reduce census? Crisis counseling? Shorter shifts?",
            "What does each cost you?",
        ],
        "footer": "Discuss, then log your decision.",
    },
    # 22: Social media
    {
        "type": "inject",
        "title": "SOCIAL MEDIA FIRESTORM",
        "subtitle": "Hour 34",
        "social_posts": [
            ("@HealthWatchFR (12,400 followers)",
             '"THREAD: My aunt is a nurse at Sainte-Claire. She says they\'ve been writing '
             'EVERYTHING by hand for 2 days. No access to patient histories. No medication '
             'barcodes. She\'s terrified of making an error. This is a patient safety EMERGENCY."'),
            ("@PatientRightsFR (34,200 followers)",
             '"If Sainte-Claire has had 180,000 patient records stolen, where is the CNIL '
             'notification? GDPR requires transparency. Silence is not acceptable. #DataBreach #GDPR"'),
            ("Dr. Laurent M. (verified physician)",
             '"The Sainte-Claire situation is exactly what we\'ve warned about. French hospitals '
             'are chronically under-invested in cybersecurity. This will happen again."'),
        ],
        "decision": "Social media is shaping the narrative. Do you engage or stay silent?",
    },
    # 23: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "A nurse's relative posted. Advocacy groups demanding answers.",
        "prompts": [
            "Official social media response? Ignore it?",
            "Ask the employee to remove it? Hire a crisis firm?",
        ],
        "footer": "Discuss, then log your decision.",
    },
    # 24: Department conflict
    {
        "type": "inject",
        "title": "INTER-DEPARTMENT CONFLICT",
        "subtitle": "Hour 36",
        "scene": {
            "label": "Scene: Incident Command Center",
            "text": ('The Chief of Surgery and the HIM Director are in a heated argument.\n\n'
                     'Chief of Surgery: "We operated for decades before computers. I am not '
                     'letting a computer problem endanger my patients by delaying their care."\n\n'
                     'HIM Director: "And we had higher error rates for decades before computers. '
                     'We are not going backward. Every mislabeled specimen is a potential catastrophe."'),
        },
        "decision": "Both are right. Both are wrong. The Incident Commander must decide.",
    },
    # 25: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "Surgery vs. HIM. Patient harm from delay vs. patient harm from error.",
        "prompts": [
            "Resume elective surgery on paper? Delay? Compromise? Governance meeting?",
        ],
        "footer": "Discuss, then log your decision.",
    },
    # 26: Day 3
    {
        "type": "critical",
        "title": "CRITICAL DECISIONS",
        "subtitle": "Wednesday -- Hour 48 to 72  |  Day 3",
        "banner": "RANSOM DEADLINE APPROACHING -- GDPR CLOCK TICKING",
        "systems": DAY3,
        "metrics": [("48h", "Downtime"), ("6,400", "Paper Records"), ("38", "Surgeries Canceled"),
                    ("\u20ac480K", "Revenue Lost"), ("89", "Complaints")],
    },
    # 27: Ransom
    {
        "type": "inject",
        "title": "THE RANSOM DEMAND",
        "subtitle": "Hour 50",
        "events": [
            ("", "The attackers have made formal contact through an encrypted channel."),
        ],
        "metrics": [("\u20ac2.5M", "Ransom"), ("\u20ac1M", "Insurance Limit"),
                    ("180K", "Records Stolen"), ("48h", "Deadline")],
        "extra_events": [
            "500 genuine patient records provided as proof.\nIncludes diagnoses, medications, national IDs, mental health and HIV status.",
            "ANSSI advises: do not pay.\nLoi LOPMI (2023): Police report within 72h of payment required for insurance.",
        ],
        "decision": "\u20ac2.5 million. 180,000 patients. Mental health records. HIV status.\nWhat do you recommend to the board?",
    },
    # 28: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "This is the hardest decision in the simulation.",
        "prompts": [
            "Pay? Negotiate? Refuse? Escalate?",
            "What are the ethical dimensions?",
            "Who bears the harm in each scenario?",
        ],
        "footer": "Discuss thoroughly, then log your decision and reasoning.",
    },
    # 29: GDPR
    {
        "type": "inject",
        "title": "GDPR 72-HOUR NOTIFICATION DEADLINE",
        "subtitle": "Hour 58 -- 14 hours remaining",
        "events": [
            ("", "GDPR Article 33: Notify CNIL within 72 hours.\nClock started 06:40 Monday. 14 hours remain."),
            ("", "Forensics:\n"
                 "  - Attackers had access 11 days before encryption\n"
                 "  - Large data transfers confirmed\n"
                 "  - Exact scope: unknown\n"
                 "  - 500-record sample verified genuine\n"
                 "  - Likely includes demographics, diagnoses, meds, labs, mental health records"),
        ],
        "decision": "File now with incomplete info? Wait for deadline? Notify patients too?",
    },
    # 30: Pause
    {
        "type": "pause",
        "title": "PAUSE -- TEAM DISCUSSION",
        "subtitle": "14 hours. Incomplete information. 180,000 potentially affected patients.",
        "prompts": [
            "File preliminary now? Wait for the deadline?",
            "How do you notify 180,000 patients?",
            "Consider: elderly patients, mental health stigma, domestic violence survivors.",
        ],
        "footer": "Discuss, then log your decision.",
    },
    # 31: Restoration
    {
        "type": "inject",
        "title": "SYSTEM RESTORATION STRATEGY",
        "subtitle": "Hour 64",
        "options": [
            ("Option A: 48-hour backup", "EHR in 12-18h. Loses pre-breach + downtime data. Clean. Back-entry: 3-4 weeks."),
            ("Option B: Decrypt & clean", "5-10 more days. Preserves data. 15-20% reinfection risk. Forensics recommends against."),
            ("Option C: New clean instance", "24-36h. Zero risk. Requires manual entry of ALL active patient data."),
            ("Option D: Hybrid", "Phase 1 (24h): clean instance for active patients.\nPhase 2 (1 week): restore backup.\nPhase 3 (2-4 weeks): merge & reconcile."),
        ],
        "decision": "Speed vs. data preservation vs. security. Who decides?",
    },
    # 32: Pause
    {
        "type": "pause",
        "title": "PAUSE -- FINAL TEAM DISCUSSION",
        "subtitle": "Four options. Each has costs. Each has risks.",
        "prompts": [
            "Which path? Technical decision? Clinical? Financial? Data integrity?",
            "Who has the final say?",
        ],
        "footer": "Discuss, then log your final decision.",
    },
    # 33: Recovery
    {
        "type": "recovery",
        "title": "RECOVERY PLANNING",
        "subtitle": "Thursday -- Hour 72+  |  Day 4",
        "body": [
            "Partial systems restoration in progress.",
        ],
        "systems": RECOVERY,
        "body2": [
            "The systems will come back in days.",
            "The trust takes months.",
            "The lessons last careers.",
        ],
        "decision": "In your student app: complete Recovery Planning.\nPrioritize back-entry. Plan trust recovery. Assess preparedness. Write your after-action summary.",
    },
    # 34: End
    {
        "type": "end",
        "title": "END OF SIMULATION",
        "body": [
            "Thank you for your work. The 2-hour debrief will follow.",
            "",
            "Export your team reports from the student app.",
            "Prepare to share your most difficult decision and what you learned.",
        ],
    },
]


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def set_slide_bg(slide, color):
    """Set the slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=18,
                          color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                          font_name="Calibri", line_spacing=1.5):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.4)
    return txBox


def add_systems_grid(slide, systems_statuses, top):
    """Add a systems status grid as text boxes."""
    left_start = Inches(0.5)
    box_w = Inches(2.4)
    box_h = Inches(0.5)
    gap_x = Inches(0.1)
    gap_y = Inches(0.08)
    cols = 5

    for i, (name, status) in enumerate(zip(SYSTEM_NAMES, systems_statuses)):
        col = i % cols
        row = i // cols
        left = left_start + col * (box_w + gap_x)
        t = top + row * (box_h + gap_y)

        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, t, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = STATUS_COLORS[status]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        label = name
        if status == "DOWN":
            label += "  [DOWN]"
        elif status == "DEGRADED":
            label += "  [DEGRADED]"
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)

    return top + (box_h + gap_y) * 2 + Inches(0.1)


def add_metrics_row(slide, metrics, top):
    """Add a row of metric boxes."""
    n = len(metrics)
    total_w = Inches(12.3)
    left_start = Inches(0.5)
    box_w = total_w / n - Inches(0.1)
    box_h = Inches(1.0)

    for i, (value, label) in enumerate(metrics):
        left = left_start + i * (box_w + Inches(0.1))
        shape = slide.shapes.add_shape(1, left, top, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = RED
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Value
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.color.rgb = RED
        p.font.bold = True
        p.font.name = "Courier New"
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)
        # Label
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(11)
        p2.font.color.rgb = MID_GRAY
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER

    return top + box_h + Inches(0.15)


def add_event_cards(slide, events, top, accent_color=RED):
    """Add event cards (timeline entries)."""
    left = Inches(0.5)
    width = Inches(12.3)
    card_h = Inches(0.6)

    for i, event_data in enumerate(events):
        if isinstance(event_data, tuple) and len(event_data) == 2:
            time_str, desc = event_data
        else:
            time_str, desc = "", str(event_data)

        text = f"{time_str}  --  {desc}" if time_str else desc
        lines = text.count("\n") + 1
        h = Inches(0.15 + lines * 0.28)

        shape = slide.shapes.add_shape(1, left, top, width, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.fill.background()

        # Left accent bar -- we simulate with a thin shape
        bar = slide.shapes.add_shape(1, left, top, Inches(0.08), h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Calibri"
        p.space_before = Pt(6)

        top += h + Inches(0.06)

    return top


def add_slide_number(slide, num):
    """Add slide number to bottom right."""
    add_textbox(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35),
                f"Slide {num + 1} / 35", font_size=10, color=MID_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_decision_prompt(slide, text, top):
    """Add a decision prompt box."""
    left = Inches(1.5)
    width = Inches(10.3)
    lines = text.count("\n") + 1
    height = Inches(0.3 + lines * 0.3)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0D, 0x2A, 0x50)
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_LIGHT
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
    return top + height + Inches(0.1)


def add_banner(slide, text, top, bg_color=RED_BANNER):
    """Add a crisis banner."""
    left = Inches(0.5)
    width = Inches(12.3)
    height = Inches(0.7)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
    return top + height + Inches(0.15)


# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for idx, slide_data in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank_layout)
        stype = slide_data["type"]

        # -- Determine background color --
        if stype == "pause":
            set_slide_bg(slide, GREEN_DARK)
        elif stype == "break":
            set_slide_bg(slide, PURPLE_BG)
        elif stype == "end":
            set_slide_bg(slide, BG_DARK)
        else:
            set_slide_bg(slide, BG_DARK)

        # -- Title --
        title_text = slide_data.get("title", "")
        title_color = RED
        if stype == "pause":
            title_color = WHITE
        elif stype == "break":
            title_color = WHITE
        elif stype == "end":
            title_color = GREEN

        add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
                    title_text, font_size=32, color=title_color, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # -- Subtitle --
        subtitle_text = slide_data.get("subtitle", "")
        if subtitle_text:
            sub_color = MID_GRAY if stype not in ("pause",) else RGBColor(0xC8, 0xE6, 0xC9)
            add_textbox(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
                        subtitle_text, font_size=16, color=sub_color,
                        alignment=PP_ALIGN.CENTER)

        current_top = Inches(1.7)

        # -- Type-specific content --
        if stype == "title":
            # Title slide
            body = slide_data.get("body", [])
            for line in body:
                line_color = LIGHT_GRAY
                line_size = 18
                line_bold = False
                if "06:12" in line:
                    line_color = RED
                    line_size = 22
                    line_bold = True
                elif "operational" in line.lower():
                    line_color = RGBColor(0x66, 0x66, 0x66)
                add_textbox(slide, Inches(1.0), current_top, Inches(11.3), Inches(0.45),
                            line, font_size=line_size, color=line_color, bold=line_bold,
                            alignment=PP_ALIGN.CENTER)
                current_top += Inches(0.45)
            current_top += Inches(0.3)
            if "systems" in slide_data:
                add_systems_grid(slide, slide_data["systems"], current_top)

        elif stype == "breach":
            events = slide_data.get("events", [])
            current_top = add_event_cards(slide, events, current_top)
            if "banner" in slide_data:
                current_top = add_banner(slide, slide_data["banner"], current_top)
            if "systems" in slide_data:
                add_systems_grid(slide, slide_data["systems"], current_top)

        elif stype == "pause":
            prompts = slide_data.get("prompts", [])
            prompt_top = Inches(2.0)
            for prompt in prompts:
                lines_count = prompt.count("\n") + 1
                h = Inches(0.2 + lines_count * 0.35)
                add_textbox(slide, Inches(1.5), prompt_top, Inches(10.3), h,
                            prompt, font_size=18, color=WHITE,
                            alignment=PP_ALIGN.LEFT)
                prompt_top += h + Inches(0.1)
            footer = slide_data.get("footer", "")
            if footer:
                add_textbox(slide, Inches(1.0), Inches(5.8), Inches(11.3), Inches(0.4),
                            footer, font_size=13, color=RGBColor(0xA5, 0xD6, 0xA7),
                            alignment=PP_ALIGN.CENTER)

        elif stype == "break":
            body = slide_data.get("body", [])
            add_multiline_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(4.0),
                                  body, font_size=20, color=PURPLE_LIGHT,
                                  alignment=PP_ALIGN.CENTER)

        elif stype in ("downtime", "escalation", "critical", "recovery"):
            if "banner" in slide_data:
                current_top = add_banner(slide, slide_data["banner"], current_top,
                                         bg_color=ORANGE_DARK if stype in ("downtime", "escalation") else RED_BANNER)
            if "systems" in slide_data:
                current_top = add_systems_grid(slide, slide_data["systems"], current_top)
            if "body" in slide_data:
                body = slide_data["body"]
                for line in body:
                    lc = LIGHT_GRAY
                    lb = False
                    ls = 18
                    if "paper" in line.lower():
                        lc = ORANGE
                        lb = True
                        ls = 22
                    add_textbox(slide, Inches(1.0), current_top, Inches(11.3), Inches(0.4),
                                line, font_size=ls, color=lc, bold=lb,
                                alignment=PP_ALIGN.CENTER)
                    current_top += Inches(0.45)
            if "metrics" in slide_data:
                current_top += Inches(0.1)
                current_top = add_metrics_row(slide, slide_data["metrics"], current_top)
            if "body2" in slide_data:
                current_top += Inches(0.2)
                for line in slide_data["body2"]:
                    lb = "days" in line or "months" in line or "careers" in line
                    add_textbox(slide, Inches(1.0), current_top, Inches(11.3), Inches(0.45),
                                line, font_size=20, color=LIGHT_GRAY, bold=lb,
                                alignment=PP_ALIGN.CENTER)
                    current_top += Inches(0.5)
            if "decision" in slide_data:
                current_top += Inches(0.1)
                add_decision_prompt(slide, slide_data["decision"], current_top)

        elif stype == "inject":
            if "events" in slide_data:
                current_top = add_event_cards(slide, slide_data["events"], current_top)
            if "metrics" in slide_data:
                current_top = add_metrics_row(slide, slide_data["metrics"], current_top)
            if "extra_events" in slide_data:
                for evt in slide_data["extra_events"]:
                    current_top = add_event_cards(slide, [("", evt)], current_top)
            if "news" in slide_data:
                news = slide_data["news"]
                # Source line
                add_textbox(slide, Inches(0.8), current_top, Inches(11.7), Inches(0.4),
                            news["source"], font_size=14, color=RED, bold=True,
                            alignment=PP_ALIGN.LEFT)
                current_top += Inches(0.5)
                # News text
                lines_count = len(news["text"]) // 100 + 1
                h = Inches(0.3 + lines_count * 0.3)
                shape = slide.shapes.add_shape(1, Inches(0.8), current_top, Inches(11.7), h)
                shape.fill.solid()
                shape.fill.fore_color.rgb = BG_CARD
                shape.line.color.rgb = RED
                shape.line.width = Pt(2)
                tf = shape.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.15)
                p = tf.paragraphs[0]
                p.text = news["text"]
                p.font.size = Pt(15)
                p.font.color.rgb = LIGHT_GRAY
                p.font.name = "Calibri"
                p.space_before = Pt(8)
                current_top += h + Inches(0.15)
            if "scene" in slide_data:
                scene = slide_data["scene"]
                add_textbox(slide, Inches(0.8), current_top, Inches(11.7), Inches(0.3),
                            scene["label"], font_size=12, color=RGBColor(0xAB, 0x47, 0xBC),
                            bold=True, alignment=PP_ALIGN.LEFT)
                current_top += Inches(0.35)
                lines_count = scene["text"].count("\n") + 1
                h = Inches(0.3 + lines_count * 0.28)
                shape = slide.shapes.add_shape(1, Inches(0.8), current_top, Inches(11.7), h)
                shape.fill.solid()
                shape.fill.fore_color.rgb = BG_CARD
                shape.line.color.rgb = RGBColor(0xAB, 0x47, 0xBC)
                shape.line.width = Pt(2)
                tf = shape.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.15)
                p = tf.paragraphs[0]
                p.text = scene["text"]
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0xE1, 0xBE, 0xE7)
                p.font.name = "Calibri"
                p.space_before = Pt(8)
                current_top += h + Inches(0.15)
            if "voicemail" in slide_data:
                vm = slide_data["voicemail"]
                add_textbox(slide, Inches(0.8), current_top, Inches(11.7), Inches(0.3),
                            vm["label"], font_size=12, color=ORANGE, bold=True,
                            alignment=PP_ALIGN.LEFT)
                current_top += Inches(0.35)
                lines_count = len(vm["text"]) // 90 + 1
                h = Inches(0.3 + lines_count * 0.25)
                shape = slide.shapes.add_shape(1, Inches(0.8), current_top, Inches(11.7), h)
                shape.fill.solid()
                shape.fill.fore_color.rgb = BG_CARD
                shape.line.color.rgb = ORANGE
                shape.line.width = Pt(2)
                tf = shape.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.15)
                p = tf.paragraphs[0]
                p.text = vm["text"]
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0xFF, 0xE0, 0xB2)
                p.font.name = "Calibri"
                p.space_before = Pt(8)
                current_top += h + Inches(0.15)
            if "social_posts" in slide_data:
                for handle, text in slide_data["social_posts"]:
                    lines_count = len(text) // 90 + 2
                    h = Inches(0.3 + lines_count * 0.25)
                    shape = slide.shapes.add_shape(1, Inches(0.8), current_top, Inches(11.7), h)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = BG_CARD
                    shape.line.color.rgb = BLUE
                    shape.line.width = Pt(2)
                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.15)
                    # Handle
                    p = tf.paragraphs[0]
                    p.text = handle
                    p.font.size = Pt(12)
                    p.font.color.rgb = BLUE
                    p.font.bold = True
                    p.font.name = "Calibri"
                    p.space_before = Pt(6)
                    # Text
                    p2 = tf.add_paragraph()
                    p2.text = text
                    p2.font.size = Pt(13)
                    p2.font.color.rgb = BLUE_LIGHT
                    p2.font.name = "Calibri"
                    current_top += h + Inches(0.08)
            if "options" in slide_data:
                for opt_title, opt_desc in slide_data["options"]:
                    lines_count = opt_desc.count("\n") + 1
                    h = Inches(0.3 + lines_count * 0.28)
                    # Determine accent color
                    accent = BLUE if "Option A" in opt_title or "Option C" in opt_title else (
                        ORANGE if "Option B" in opt_title else RGBColor(0xAB, 0x47, 0xBC))
                    shape = slide.shapes.add_shape(1, Inches(0.8), current_top, Inches(11.7), h)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = BG_CARD
                    shape.line.fill.background()
                    # Accent bar
                    bar = slide.shapes.add_shape(1, Inches(0.8), current_top, Inches(0.08), h)
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = accent
                    bar.line.fill.background()

                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.2)
                    p = tf.paragraphs[0]
                    p.text = f"{opt_title}: {opt_desc}"
                    p.font.size = Pt(13)
                    p.font.color.rgb = LIGHT_GRAY
                    p.font.name = "Calibri"
                    p.space_before = Pt(6)
                    current_top += h + Inches(0.06)

            if "decision" in slide_data:
                current_top += Inches(0.05)
                add_decision_prompt(slide, slide_data["decision"], current_top)

        elif stype == "end":
            body = slide_data.get("body", [])
            add_multiline_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(3.0),
                                  body, font_size=18, color=MID_GRAY,
                                  alignment=PP_ALIGN.CENTER)

        # -- Slide number --
        add_slide_number(slide, idx)

        # -- Speaker notes --
        narration = NARRATIONS.get(idx, "")
        if narration:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = narration

    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "SHIFT-Cybersecurity-Simulation.pptx")
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
